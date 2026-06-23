"""
Générateurs d'export PPT et Excel pour les livrables M&A.

Sell-side : PPT au format Inspirit Partners (charte du template 3b)
Buy-side  : PPT cibles de croissance + Excel screening
"""
from __future__ import annotations
import io, re, datetime, textwrap
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import copy

# ── Charte couleurs (extraites du template 3b) ────────────────────────────────
TEAL       = RGBColor(0x00, 0x87, 0x8E)
DARK_TEAL  = RGBColor(0x00, 0x6B, 0x70)
DARK_NAVY  = RGBColor(0x1A, 0x3A, 0x3C)
ORANGE     = RGBColor(0xF1, 0x59, 0x22)
CARD_GRAY  = RGBColor(0xE8, 0xEB, 0xED)
MID_GRAY   = RGBColor(0x6B, 0x72, 0x80)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLACK      = RGBColor(0x00, 0x00, 0x00)
BG_GRAY    = RGBColor(0xF5, 0xF7, 0xF8)

# Buy-side (charte app)
BUY_NAVY   = RGBColor(0x1A, 0x27, 0x44)
BUY_GREEN  = RGBColor(0x2E, 0x7D, 0x55)
BUY_LIGHT  = RGBColor(0xE8, 0xF5, 0xEE)

# ── Dimensions slide (A4 paysage, comme le template) ─────────────────────────
SLIDE_W = Emu(10691813)   # 11.693 in
SLIDE_H = Emu(7559675)    # 8.267 in


# ── Helpers ───────────────────────────────────────────────────────────────────

def _rgb_elem(rgb: RGBColor):
    e = etree.Element("{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr")
    e.set("val", f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}")
    return e


def _add_rect(slide, l, t, w, h, rgb: RGBColor, alpha=None):
    """Ajoute un rectangle coloré plein."""
    sp = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    sp.line.fill.background()
    fill = sp.fill
    fill.solid()
    fill.fore_color.rgb = rgb
    return sp


def _add_textbox(slide, l, t, w, h, text, size_pt, bold=False, color=BLACK,
                 align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def _set_bg(slide, rgb: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _parse_sections(content: str) -> list[dict]:
    """
    Découpe le markdown en sections.
    Chaque ## ou ### devient une slide.
    Retourne : [{title, teaser, body_lines}]
    """
    sections = []
    current = None

    for line in content.splitlines():
        if line.startswith("## ") or line.startswith("### "):
            if current:
                sections.append(current)
            title = line.lstrip("# ").strip()
            current = {"title": title, "teaser": "", "body_lines": []}
        elif current is not None:
            stripped = line.strip()
            if not current["teaser"] and stripped and not stripped.startswith("|") and not stripped.startswith("*"):
                current["teaser"] = stripped
            else:
                current["body_lines"].append(line)

    if current:
        sections.append(current)

    # Si aucune section trouvée, mettre tout dans une slide
    if not sections:
        lines = [l for l in content.splitlines() if l.strip()]
        teaser = lines[0] if lines else ""
        body = "\n".join(lines[1:]) if len(lines) > 1 else ""
        sections = [{"title": "Contenu", "teaser": teaser, "body_lines": body.splitlines()}]

    return sections


def _section_body_text(section: dict, max_chars=900) -> str:
    """Renvoie le texte body nettoyé, tronqué si besoin."""
    lines = []
    for line in section["body_lines"]:
        s = line.strip()
        if not s:
            continue
        # Nettoyer le markdown
        s = re.sub(r"\*\*([^*]+)\*\*", r"\1", s)   # **bold**
        s = re.sub(r"\*([^*]+)\*", r"\1", s)         # *italic*
        s = re.sub(r"^[-•*]\s*", "• ", s)            # bullets
        s = re.sub(r"^\d+\.\s*", "→ ", s)            # numbered
        s = re.sub(r"^\|.*\|$", "", s)               # tables (skip)
        s = re.sub(r"^#+\s*", "", s)                  # nested headers
        if s:
            lines.append(s)
    text = "\n".join(lines)
    if len(text) > max_chars:
        text = text[:max_chars].rsplit("\n", 1)[0] + "\n…"
    return text


# ── SELL-SIDE PPT ─────────────────────────────────────────────────────────────

def _sell_slide(prs: Presentation, section: dict, company: str, section_num: int, total: int, date_str: str):
    """Ajoute une slide sell-side au format Inspirit Partners."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    _set_bg(slide, BG_GRAY)

    W, H = SLIDE_W, SLIDE_H
    LEFT_COL  = Cm(0.5)
    HEADER_H  = Cm(1.0)
    MARGIN    = Cm(0.5)

    # ── Barre header ──────────────────────────────────────────────────────────
    _add_rect(slide, 0, 0, W, HEADER_H, DARK_NAVY)
    hdr_text = f"CONFIDENTIEL  ·  {company}  ·  Mémorandum d'information  ·  {date_str}"
    _add_textbox(slide, MARGIN, Cm(0.15), W - Cm(1), HEADER_H - Cm(0.3),
                 hdr_text, 7.5, False, WHITE)

    # Numéro de slide en haut à droite
    _add_textbox(slide, W - Cm(1.5), Cm(0.15), Cm(1.2), HEADER_H - Cm(0.3),
                 f"{section_num}/{total}", 7.5, False, RGBColor(0x9C, 0xB4, 0xB6),
                 align=PP_ALIGN.RIGHT)

    # ── Bande teal sous le header ─────────────────────────────────────────────
    _add_rect(slide, 0, HEADER_H, W, Cm(0.07), TEAL)

    # ── Titre de la slide ─────────────────────────────────────────────────────
    title_text = section["title"].upper() if len(section["title"]) < 60 else section["title"]
    txTitle = _add_textbox(slide, LEFT_COL, Cm(1.25), W - Cm(1), Cm(0.8),
                           title_text, 15, True, TEAL)

    # ── Ligne séparatrice sous le titre ───────────────────────────────────────
    _add_rect(slide, LEFT_COL, Cm(2.1), Cm(8), Cm(0.04), TEAL)

    # ── Teaser (accroche italique) ────────────────────────────────────────────
    teaser = section.get("teaser", "")
    if teaser:
        _add_textbox(slide, LEFT_COL, Cm(2.25), W - Cm(1), Cm(0.9),
                     teaser, 9.5, False, RGBColor(0x2C, 0x3E, 0x50), italic=True)

    # ── Carte de contenu ──────────────────────────────────────────────────────
    card_top = Cm(3.2) if teaser else Cm(2.5)
    card_h = H - card_top - Cm(0.6)

    # Fond de la carte (blanc légèrement grisé)
    card_rect = _add_rect(slide, LEFT_COL, card_top, W - Cm(1), card_h, WHITE)
    # Bordure gauche teal
    _add_rect(slide, LEFT_COL, card_top, Cm(0.12), card_h, TEAL)

    # Contenu
    body_text = _section_body_text(section, max_chars=1100)
    if body_text:
        _add_textbox(slide, LEFT_COL + Cm(0.3), card_top + Cm(0.2),
                     W - Cm(1.5), card_h - Cm(0.4),
                     body_text, 9, False, BLACK)

    # ── Étiquette de section (coin bas gauche) ────────────────────────────────
    _add_rect(slide, 0, H - Cm(0.5), Cm(5), Cm(0.5), DARK_NAVY)
    _add_textbox(slide, MARGIN, H - Cm(0.5), Cm(4.5), Cm(0.5),
                 company.upper(), 7, True, RGBColor(0x9C, 0xB4, 0xB6))

    return slide


def generate_sell_pptx(content: str, company: str) -> bytes:
    """Génère un fichier PPTX sell-side à partir du texte Markdown de l'IA."""
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    sections = _parse_sections(content)
    date_str = datetime.date.today().strftime("%B %Y")

    for i, section in enumerate(sections, 1):
        _sell_slide(prs, section, company, i, len(sections), date_str)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ── BUY-SIDE PPT ──────────────────────────────────────────────────────────────

def _buy_slide(prs: Presentation, section: dict, company: str, slide_num: int,
               total: int, slide_type: str = "content"):
    """Ajoute une slide buy-side (Cibles de croissance)."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    _set_bg(slide, BG_GRAY)

    W, H = SLIDE_W, SLIDE_H
    MARGIN = Cm(0.6)
    HEADER_H = Cm(0.95)

    # ── Header ────────────────────────────────────────────────────────────────
    _add_rect(slide, 0, 0, W, HEADER_H, BUY_NAVY)
    hdr = f"Cibles de croissance externe  ·  {company}  ·  CONFIDENTIEL"
    _add_textbox(slide, MARGIN, Cm(0.18), W - Cm(2), Cm(0.65), hdr, 7.5, False, WHITE)
    _add_textbox(slide, W - Cm(1.8), Cm(0.18), Cm(1.5), Cm(0.65),
                 f"{slide_num}/{total}", 7.5, False, RGBColor(0x9C, 0xB4, 0xB6),
                 align=PP_ALIGN.RIGHT)

    # Bande verte sous header
    _add_rect(slide, 0, HEADER_H, W, Cm(0.06), BUY_GREEN)

    # ── Sidebar gauche verte ──────────────────────────────────────────────────
    SIDEBAR_W = Cm(0.4)
    _add_rect(slide, 0, HEADER_H + Cm(0.06), SIDEBAR_W, H - HEADER_H - Cm(0.06), BUY_GREEN)

    # ── Titre ─────────────────────────────────────────────────────────────────
    title = section["title"]
    _add_textbox(slide, MARGIN + SIDEBAR_W, Cm(1.15), W - MARGIN - SIDEBAR_W - Cm(0.4),
                 Cm(0.75), title, 14, True, BUY_NAVY)

    _add_rect(slide, MARGIN + SIDEBAR_W, Cm(2.0), Cm(7), Cm(0.04), BUY_GREEN)

    # ── Teaser ────────────────────────────────────────────────────────────────
    teaser = section.get("teaser", "")
    card_top = Cm(2.1)
    if teaser:
        _add_textbox(slide, MARGIN + SIDEBAR_W, Cm(2.15), W - MARGIN * 2 - SIDEBAR_W, Cm(0.85),
                     teaser, 9, False, RGBColor(0x2C, 0x3E, 0x50), italic=True)
        card_top = Cm(3.05)

    # ── Carte contenu ─────────────────────────────────────────────────────────
    card_h = H - card_top - Cm(0.55)
    _add_rect(slide, MARGIN + SIDEBAR_W, card_top, W - MARGIN * 2 - SIDEBAR_W, card_h, WHITE)
    _add_rect(slide, MARGIN + SIDEBAR_W, card_top, Cm(0.1), card_h, BUY_GREEN)

    body = _section_body_text(section, max_chars=1200)
    if body:
        _add_textbox(slide, MARGIN + SIDEBAR_W + Cm(0.25), card_top + Cm(0.2),
                     W - MARGIN * 2 - SIDEBAR_W - Cm(0.5), card_h - Cm(0.4),
                     body, 9, False, BLACK)

    # ── Footer ────────────────────────────────────────────────────────────────
    _add_rect(slide, 0, H - Cm(0.45), W, Cm(0.45), BUY_NAVY)
    _add_textbox(slide, MARGIN, H - Cm(0.45), Cm(6), Cm(0.45),
                 company.upper(), 7, True, RGBColor(0x9C, 0xB4, 0xB6))

    return slide


def generate_buy_pptx(content: str, company: str, module_title: str = "Cibles de croissance externe") -> bytes:
    """Génère un fichier PPTX buy-side à partir du texte Markdown de l'IA."""
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    sections = _parse_sections(content)
    date_str = datetime.date.today().strftime("%B %Y")

    # Slide de titre
    slide_layout = prs.slide_layouts[6]
    title_slide = prs.slides.add_slide(slide_layout)
    _set_bg(title_slide, BUY_NAVY)
    W, H = SLIDE_W, SLIDE_H
    _add_textbox(title_slide, Cm(1), Cm(2.5), W - Cm(2), Cm(2),
                 company, 32, True, WHITE, align=PP_ALIGN.CENTER)
    _add_rect(title_slide, W / 4, Cm(4.6), W / 2, Cm(0.08), BUY_GREEN)
    _add_textbox(title_slide, Cm(1), Cm(4.9), W - Cm(2), Cm(1.2),
                 module_title, 20, False, RGBColor(0x9C, 0xB4, 0xB6), align=PP_ALIGN.CENTER)
    _add_textbox(title_slide, Cm(1), H - Cm(1.5), W - Cm(2), Cm(1),
                 date_str + "  ·  CONFIDENTIEL", 9, False,
                 RGBColor(0x6B, 0x82, 0x99), align=PP_ALIGN.CENTER)

    for i, section in enumerate(sections, 2):
        _buy_slide(prs, section, company, i, len(sections) + 1)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ── EXCEL SCREENING ───────────────────────────────────────────────────────────

def generate_screening_xlsx(content: str, company: str) -> bytes:
    """
    Génère un Excel de screening M&A à partir du markdown.
    Colonnes : Société, Activité, Chiffres clés, Actionnariat, Opérations M&A,
               Sources, Pertinence offre, Pertinence clientèle, Conclusion
    """
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side, numbers
    from openpyxl.utils import get_column_letter

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Screening"

    NAVY_HEX  = "1A2744"
    GREEN_HEX = "2E7D55"
    LIGHT_HEX = "E8F5EE"
    GRAY_HEX  = "F5F7F8"

    thin = Side(style="thin", color="CCCCCC")
    border = Border(left=thin, right=thin, top=thin, bottom=thin)

    COLS = [
        ("#", 5),
        ("Société", 22),
        ("Activité principale", 32),
        ("Chiffres clés", 22),
        ("Actionnariat", 20),
        ("Opérations M&A", 24),
        ("Pertinence offre", 18),
        ("Pertinence clientèle", 18),
        ("Conclusion", 16),
        ("Sources", 28),
    ]

    # ── Ligne de titre ────────────────────────────────────────────────────────
    ws.merge_cells("A1:J1")
    title_cell = ws["A1"]
    title_cell.value = f"Screening M&A — {company}"
    title_cell.font = Font(name="Arial", size=14, bold=True, color="FFFFFF")
    title_cell.fill = PatternFill("solid", fgColor=NAVY_HEX)
    title_cell.alignment = Alignment(horizontal="left", vertical="center")
    ws.row_dimensions[1].height = 28

    # ── En-têtes ──────────────────────────────────────────────────────────────
    for col_idx, (col_name, col_w) in enumerate(COLS, 1):
        cell = ws.cell(row=2, column=col_idx, value=col_name)
        cell.font = Font(name="Arial", size=9, bold=True, color="FFFFFF")
        cell.fill = PatternFill("solid", fgColor=GREEN_HEX)
        cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
        cell.border = border
        ws.column_dimensions[get_column_letter(col_idx)].width = col_w
    ws.row_dimensions[2].height = 30

    # ── Parser le contenu ─────────────────────────────────────────────────────
    companies = _parse_companies_from_md(content)

    for row_idx, co in enumerate(companies, 3):
        fill_color = LIGHT_HEX if row_idx % 2 == 1 else "FFFFFF"
        values = [
            row_idx - 2,
            co.get("nom", "NC"),
            co.get("activite", "NC"),
            co.get("chiffres_cles", "NC"),
            co.get("actionnariat", "NC"),
            co.get("manda", "NC"),
            co.get("pertinence_offre", "NC"),
            co.get("pertinence_clientele", "NC"),
            co.get("conclusion", "NC"),
            co.get("sources", "NC"),
        ]
        for col_idx, val in enumerate(values, 1):
            cell = ws.cell(row=row_idx, column=col_idx, value=val)
            cell.font = Font(name="Arial", size=8.5)
            cell.fill = PatternFill("solid", fgColor=fill_color)
            cell.alignment = Alignment(vertical="top", wrap_text=True)
            cell.border = border
        ws.row_dimensions[row_idx].height = 45

    # ── Freeze panes ──────────────────────────────────────────────────────────
    ws.freeze_panes = "B3"

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _parse_companies_from_md(content: str) -> list[dict]:
    """
    Extrait les fiches entreprises du markdown généré par buy_04 / buy_05.
    Cherche les blocs ### Concurrent N° — NOM ou ### [N°. Nom]
    """
    companies = []

    # Séparer par blocs
    blocks = re.split(r"\n(?=###\s+(?:Concurrent|Cible|\d+[\.\)]|[A-Z]))", content)

    for block in blocks:
        if not block.strip():
            continue
        co: dict = {}

        # Nom
        header = block.splitlines()[0] if block.splitlines() else ""
        name_m = re.search(r"###\s+(?:Concurrent\s+\d+\s*[—\-:]+\s*)?(.+)", header)
        if name_m:
            co["nom"] = name_m.group(1).strip().lstrip("*").rstrip("*")

        # Extraire les champs des tableaux Markdown
        for line in block.splitlines():
            m = re.match(r"\|\s*([^|]+?)\s*\|\s*([^|]+?)\s*\|", line)
            if not m:
                continue
            key = m.group(1).strip().lower()
            val = m.group(2).strip()
            if val in ("---", "Information", "Lien", "Champ", "---", ""):
                continue

            if any(k in key for k in ["activité principale", "activite principale", "activité"]):
                co.setdefault("activite", val)
            elif any(k in key for k in ["ca estimé", "ca ", "chiffre"]):
                co.setdefault("chiffres_cles", val)
            elif "effectif" in key:
                existing = co.get("chiffres_cles", "")
                co["chiffres_cles"] = (existing + f"\nEffectifs : {val}").strip()
            elif "actionnariat" in key or "statut" in key:
                co.setdefault("actionnariat", val)
            elif "opération" in key or "m&a" in key:
                co.setdefault("manda", val)
            elif "point différenciant" in key or "différenciant" in key:
                co.setdefault("pertinence_offre", val)
            elif "source" in key:
                co.setdefault("sources", val)

        if co.get("nom"):
            companies.append(co)

    return companies if companies else [{"nom": "À compléter", "activite": content[:200]}]


# ── Exécution du code python-pptx généré par l'IA (sell_03) ──────────────────

_DANGEROUS = [
    "subprocess", "os.system", "os.popen", "__import__",
    "eval(", "exec(", "open(", "shutil", "socket",
]

def try_exec_pptx_code(ai_output: str) -> bytes | None:
    """
    Si la réponse de l'IA contient un bloc ```python``` avec du code python-pptx,
    l'exécute et renvoie les bytes du PPTX résultant.
    Renvoie None si pas de code ou si l'exécution échoue.
    """
    # Extraire le ou les blocs python
    blocks = re.findall(r"```python\s*(.*?)```", ai_output, re.DOTALL)
    if not blocks:
        # Essayer aussi les blocs sans langage
        blocks = re.findall(r"```\s*(from pptx|import pptx.*?)```", ai_output, re.DOTALL)
    if not blocks:
        return None

    code = "\n".join(blocks)

    # Vérification de sécurité minimale
    for danger in _DANGEROUS:
        if danger in code:
            return None

    # Remplacer prs.save("quelquechose") par prs.save(_output_buf)
    code = re.sub(r'prs\.save\(["\'][^"\']*["\']\)', "prs.save(_output_buf)", code)
    # Même chose avec des variables : prs.save(output_path) etc.
    code = re.sub(r'prs\.save\(\w+\)', "prs.save(_output_buf)", code)

    buf = io.BytesIO()
    namespace: dict = {"_output_buf": buf}

    try:
        exec(textwrap.dedent(code), namespace)  # noqa: S102
        buf.seek(0)
        data = buf.read()
        if len(data) < 1000:          # trop petit = probablement pas un vrai PPTX
            return None
        return data
    except Exception:
        return None
